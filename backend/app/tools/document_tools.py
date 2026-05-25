from __future__ import annotations

import csv
from pathlib import Path
from typing import Any

from app.core.paths import resolve_authorized
from app.indexer.ocr_service import extract_pdf_text_with_ocr_fallback
from app.policy.risk import RiskLevel
from app.services import document_service
from app.tools.schemas import ToolDefinition


_EXTRACT_TEXT_LIMIT = 20000
_CHUNK_CHARS = document_service.DEFAULT_CHUNK_CHARS


def _allowed(context: dict[str, Any]) -> list[str]:
    return list(context.get("allowed_directories") or [])


def _document_max_chars_to_llm(context: dict[str, Any]) -> int:
    settings = context.get("settings")
    value = getattr(settings, "document_max_chars_to_llm", document_service.DEFAULT_MAX_CHARS_TO_LLM)
    return max(1, int(value))


def _provider(task: str = "subagent"):
    return document_service._provider(task)


def _chunk_text(text: str, chunk_chars: int = _CHUNK_CHARS) -> list[str]:
    return document_service.chunk_document(
        text,
        chunk_chars=chunk_chars,
        overlap=0,
        max_chunks=None,
        max_chars=max(len(text or ""), 1),
    )


def _rank_chunks(query: str, chunks: list[str]) -> list[str]:
    return [chunk.text for chunk in document_service.rank_chunks(query, chunks, top_k=len(chunks) or 1)]


def extract_text_from_path(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in {".txt", ".md", ".json", ".csv", ".py", ".ts", ".tsx", ".js", ".yaml", ".yml"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".pdf":
        return extract_pdf_text_with_ocr_fallback(path)
    if ext == ".docx":
        try:
            from docx import Document

            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as exc:
            return f"[DOCX extraction unavailable: {exc}]"
    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook

            wb = load_workbook(path, read_only=True, data_only=True)
            lines: list[str] = []
            for ws in wb.worksheets:
                lines.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    lines.append(",".join("" if value is None else str(value) for value in row))
            return "\n".join(lines)
        except Exception as exc:
            return f"[XLSX extraction unavailable: {exc}]"
    if ext == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            lines = []
            for idx, slide in enumerate(prs.slides, start=1):
                lines.append(f"# Slide {idx}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            return "\n".join(lines)
        except Exception as exc:
            return f"[PPTX extraction unavailable: {exc}]"
    return "[Unsupported document type]"


def extract_text(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    return {
        "path": str(path),
        "text": text[:_EXTRACT_TEXT_LIMIT],
        "truncated": len(text) > _EXTRACT_TEXT_LIMIT,
    }


def summarize(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    result = document_service.summarize_text(
        text,
        path_label=path.name,
        max_chars_to_llm=_document_max_chars_to_llm(context),
        provider_resolver=_provider,
    )
    return {"path": str(path), **result}


def qa(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    question = str(args.get("question") or "").strip()
    text = extract_text_from_path(path)
    result = document_service.answer_question(
        text,
        question,
        path_label=path.name,
        max_chars_to_llm=_document_max_chars_to_llm(context),
        provider_resolver=_provider,
    )
    return {"path": str(path), **result}


def convert_to_markdown(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    return {"markdown": f"# {path.name}\n\n{text}"[:_EXTRACT_TEXT_LIMIT]}


def analyze_csv(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    rows = list(csv.DictReader(path.open("r", encoding="utf-8", errors="ignore")))
    return {"path": str(path), "rows": len(rows), "columns": list(rows[0].keys()) if rows else []}


def analyze_xlsx(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    return {"path": str(path), "preview": text[:2000]}


def generate_report(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    content = str(args.get("content") or "").strip()
    title = str(args.get("title") or "Report").strip() or "Report"
    return document_service.generate_report(
        content,
        title=title,
        max_chars_to_llm=_document_max_chars_to_llm(context),
        provider_resolver=_provider,
    )


def register(registry) -> None:
    defs = [
        ("document.extract_text", extract_text),
        ("document.summarize", summarize),
        ("document.qa", qa),
        ("document.convert_to_markdown", convert_to_markdown),
        ("document.analyze_csv", analyze_csv),
        ("document.analyze_xlsx", analyze_xlsx),
        ("document.generate_report", generate_report),
    ]
    for name, fn in defs:
        registry.register(
            ToolDefinition(
                name=name,
                description=name.replace(".", " "),
                input_schema={},
                output_schema={},
                risk_level=RiskLevel.R0_READ_ONLY,
                agent_owner="DocumentAgent",
                supports_dry_run=False,
                requires_authorized_path=name != "document.generate_report",
                execute=fn,
            )
        )
